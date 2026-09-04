# -*- coding: utf-8 -*-
"""Объединённая колода: каркас Яна (структура/тема) + детали из промежуточного отчёта."""

import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

SRC, OUT = 'source.pptx', 'merged.pptx'


def _norm(s):
    return s.replace('\x0b', '').replace('\r', '').strip()


def replace_in_paragraph(p, new_text):
    runs = p.runs
    if not runs:
        p.add_run().text = new_text
        return
    runs[0].text = new_text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def apply_replacements(tf, mapping):
    m = {_norm(k): v for k, v in mapping.items()}
    full = '\n'.join(p.text for p in tf.paragraphs)
    if _norm(full) in m:
        parts = m[_norm(full)].split('\n')
        for i, p in enumerate(tf.paragraphs):
            replace_in_paragraph(p, parts[i] if i < len(parts) else '')
        return True
    hit = False
    for p in tf.paragraphs:
        if _norm(p.text) in m:
            replace_in_paragraph(p, m[_norm(p.text)])
            hit = True
    return hit


def apply_to_slide(slide, mapping):
    used = set()
    for s in slide.shapes:
        if s.has_text_frame:
            before = set(m.values())
            if apply_replacements(s.text_frame, mapping):
                used |= {k for k in mapping if _norm(k) not in
                         {c.text for sh in slide.shapes
                          if sh.has_text_frame for pp in sh.text_frame.paragraphs
                          for c in [pp]} if k in mapping} if False else set()
    return used


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


prs = Presentation(SRC)

# ---- 1. клон слайда 5 (stat-находка) -> НАХОДКА 04 (матрица ширины) ----
src = prs.slides[4]
new = prs.slides.add_slide(src.slide_layout)
for shp in list(new.shapes):
    shp._element.getparent().remove(shp._element)
src_csld = src._element.find(qn('p:cSld'))
new_csld = new._element.find(qn('p:cSld'))
src_bg = src_csld.find(qn('p:bg'))
if src_bg is not None:
    new_csld.insert(0, copy.deepcopy(src_bg))
for shp in src.shapes:
    new.shapes._spTree.append(copy.deepcopy(shp._element))
# переставить: клон был добавлен в конец -> вставить после слайда 7 (idx 6)
lst = prs.slides._sldIdLst
ids = list(lst)
lst.remove(ids[-1])
lst.insert(7, ids[-1])

# ---- 2. замены по слайдам (scope: индекс слайда) ----
R = {}

R[2] = {  # слайд 3 РЕШЕНИЕ: свежие цифры memred
    '22 шаблона (21 класс) + 3 APT-цепочки. MSI 0–100, SARIF 2.1.0, findings.json. Ширина и готовность к ASOC.':
        '24 шаблона (23 класса) + 3 APT-цепочки. MSI 0–100, SARIF 2.1.0, proof-отчёт. 19 подтверждений из 57 прогонов.',
}
R[3] = {  # слайд 4 АРХИТЕКТУРА: proof-цепочка в evidence-слой
    'Чтение состояния памяти (Mongo / GET /memory) + трасса Langfuse → диф «до/после». Не голый ответ модели.':
        'Цепочка: память ДО → ПОСЛЕ → оценка → итог. sha256 артефактов, --verify ловит правки.',
}
R[4] = {  # слайд 5 НАХОДКА 01: наши подтверждения
    'Личная просьба клиента оседает в agent_policy_memories без user_id → становится «правилом для всех клиентов» и переживает сброс сессии.':
        'Личная просьба оседает в agent_policy_memories без user_id → «правило для всех». У нас 8/8 атак: тайны, реквизиты, переводы.',
}
R[7] = {  # клон -> НАХОДКА 04: матрица ширины
    'НАХОДКА 01  ·  ПРИОРИТЕТ №1': 'НАХОДКА 04  ·  МАТРИЦА ШИРИНЫ',
    'Broken Access Control через\nглобальный слой памяти':
        'Запись яда не зависит\nот стиля формулировки',
    'Личная просьба клиента оседает в agent_policy_memories без user_id → становится «правилом для всех клиентов» и переживает сброс сессии.':
        'Одна цель — 16 формулировок: стили, языки, форматы, роли. Все 16 записались в память. Активировалась одна: узкое место — не стиль, а recall.',
    '100%': '16/16',
    'запись в global\nслой (strong/weak)': 'формулировок одной цели\nзаписались в память',
    '2/3': '1/16',
    'APT CH-1:\nполное досье чужого клиента': 'активировались\nв новой сессии',
    'BAC': 'ИТОГ',
    'утечка ФИО + счёт\n+ состав портфеля': 'фильтр семантики\nзаписи, не стиля',
    'Подтверждено обоими движками. Приоритет №1 кейсодателя.':
        'Стилистические детекторы записи не сработают — нужен смысловой контроль.',
    '05': '08',
}
R[8] = {  # старый 08 ДОКАЗАТЕЛЬСТВО
    'Снапшот до/после. Что именно записалось в agent_policy_memories / user memory.':
        'Цепочка: атака → память ДО → ПОСЛЕ → оценка → итог. Каждый прогон подписан sha256.',
    'Вспомогательный. Мягче, даёт ложноположительные. Согласие двух вердикторов — 89% (κ=0.73).':
        'Вспомогательный. Мягче, даёт ложноположительные. Согласие двух вердикторов — 88% (κ=0.71).',
    '08': '09',
}
R[9] = {  # старый 09 ПОКРЫТИЕ
    '55 находок в последнем прогоне memred: 10 critical / 11 high / 32 medium':
        '57 находок: 10 critical / 11 high / 34 medium · 111 триггеров',
    '09': '10',
}
R[10] = {  # старый 10 MVP
    '●  Evidence-слой (Mongo / Langfuse)': '●  Отчёт-доказательство (sha256 + verify)',
    '10': '11',
}
R[11] = {'11': '12'}
R[12] = {'12': '13'}

for idx, mapping in R.items():
    slide = prs.slides[idx]
    for s in slide.shapes:
        if s.has_text_frame:
            apply_replacements(s.text_frame, mapping)

# ---- 3. заметки докладчика (лёгкий нарратив) ----
NOTES = {
    0: 'Открываем: мы ломаем память агента до выката в прод. Два независимых движка, 30+ классов атак, приоритет заказчика — Broken Access Control.',
    1: 'Проблема: память — это не «запрос-ответ». Отравленный факт живёт после сброса сессии и всплывает у другого клиента. Существующие сканеры это не видят.',
    2: 'Наш цикл: записали яд → активировали в новой сессии → показали последствие → собрали доказательства. Два движка проверяют друг друга: core-harness и memred.',
    3: 'Архитектура: атака — это файл, вход — API-контракт, доказательства — из памяти, а не из ответа модели. ИБ может дополнить батарею своими атаками без кода.',
    4: 'Главная находка: личная просьба клиента становится правилом для всех. Записали тайну, реквизиты, инструкцию — всё осело в глобальном слое памяти.',
    5: 'Режим protected — это авторизация чтения. Запись яда он не блокирует: write-path открыт в обоих режимах. Вывод: защищать нужно запись.',
    6: 'Цепочки воспроизводятся: флуд и фейковое «мы договаривались» активируются стабильно, до 9 из 9. На этом можно строить go/no-go.',
    7: 'Эксперимент ширины: одна цель, 16 формулировок — все 16 дошли в память, активировалась одна. Значит, детектор должен смотреть на смысл записываемого, а не на стиль.',
    8: 'Доказательство: по каждой атаке — цепочка «атака → память до → память после → оценка → итог», прогоны подписаны хэшами. Судья LLM вспомогательный: согласие 88%.',
    9: 'Покрытие: 57 находок, 10 критических. Все три сценария кейсодателя закрыты: cross-user, persistent poisoning, tool-calling. Выгрузка в ASOC готова.',
    10: 'Честные границы: внутри — CLI, black-box, evidence, pluggable-атаки. Вне — jailbreak-и и обещание «полной защищённости».',
    11: 'Как пользуется ИБ: задал контракт → запустил батарею → прочитал finding → принял go/no-go. Встраивается в MLSecOps, а не остаётся демо.',
    12: 'Финал: память — поверхность атаки, protected её не закрывает. Мы даём доказательства и расширяемость. Команда: Яна, Артур, Андрей.',
}
for idx, text in NOTES.items():
    set_notes(prs.slides[idx], text)

prs.save(OUT)

# ---- 4. проверка ----
chk = Presentation(OUT)
print('slides:', len(chk.slides._sldIdLst))
for i, s in enumerate(chk.slides, 1):
    texts = [sh.text_frame.text for sh in s.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
    joined = ' | '.join(t.replace('\n', ' ')[:40] for t in texts[:3])
    print(f'{i:>2}: {joined[:100]}')
